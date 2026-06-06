from pathlib import Path
from pptx import Presentation

source = Path(r"C:\Users\C0SU\Downloads\Overview_Organisational_Behaviour_CRKC_5001 (2).pptx")
output = Path("tmp/pptx-extract")
media = output / "media"
media.mkdir(parents=True, exist_ok=True)

prs = Presentation(source)
lines = [f"# PPTX Extract\n\nSlides: {len(prs.slides)}\n"]
for index, slide in enumerate(prs.slides, start=1):
    lines.append(f"\n## Slide {index}\n")
    for shape_index, shape in enumerate(slide.shapes, start=1):
        if hasattr(shape, "text") and shape.text.strip():
            lines.append(shape.text.strip() + "\n")
        if shape.shape_type == 13:
            ext = shape.image.ext
            target = media / f"slide-{index:02d}-image-{shape_index:02d}.{ext}"
            target.write_bytes(shape.image.blob)
            lines.append(f"[Image: {target.as_posix()}]\n")

(output / "extract.md").write_text("\n".join(lines), encoding="utf-8")
print(f"slides={len(prs.slides)} media={len(list(media.iterdir()))}")
